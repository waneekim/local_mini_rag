#!/usr/bin/env python3
import argparse
import csv
import importlib.util
import io
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".html", ".htm", ".log"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"}


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--self-check", action="store_true")
    parser.add_argument("--glossary", action="store_true")
    args = parser.parse_args()

    if args.self_check:
        print(json.dumps(self_check(), ensure_ascii=False, indent=2))
        return

    request = json.loads(sys.stdin.read() or "{}")
    result = extract_glossary_rows(request) if args.glossary else extract(request)
    print(json.dumps(result, ensure_ascii=False))


def self_check():
    return {
        "python": sys.version.split()[0],
        "optional_modules": {
            "fitz": has_module("fitz"),
            "docx": has_module("docx"),
            "pptx": has_module("pptx"),
            "openpyxl": has_module("openpyxl"),
            "PIL": has_module("PIL"),
            "pytesseract": has_module("pytesseract"),
        },
        "binaries": {
            "tesseract": shutil.which("tesseract") or "",
            "soffice": shutil.which("soffice") or "",
        },
    }


def extract(request):
    try:
        kind = request.get("kind") or ""
        file_path = request.get("filePath") or ""
        text = request.get("text") or ""
        relative_path = request.get("relativePath") or request.get("fileName") or ""
        title = request.get("title") or request.get("fileName") or "source"
        suffix = Path(file_path or relative_path or title).suffix.lower()

        if kind == "text" or text:
            return ok([unit(text, request, {"input": "text"})])
        if suffix in TEXT_EXTENSIONS:
            return ok(extract_text_file(file_path, request))
        if suffix == ".pdf":
            return ok(extract_pdf(file_path, request))
        if suffix == ".docx":
            return ok(extract_docx(file_path, request))
        if suffix == ".doc":
            return ok(extract_legacy_office(file_path, request, ".docx"))
        if suffix == ".pptx":
            return ok(extract_pptx(file_path, request))
        if suffix == ".ppt":
            return ok(extract_legacy_office(file_path, request, ".pptx"))
        if suffix in {".xlsx", ".xlsm"}:
            return ok(extract_xlsx(file_path, request))
        if suffix == ".xls":
            return ok(extract_legacy_office(file_path, request, ".xlsx"))
        if suffix in IMAGE_EXTENSIONS:
            return ok(extract_image(file_path, request))
        return fail(f"Unsupported file type '{suffix or kind}'. Convert it to PDF, text, docx, pptx, xlsx, or image first.")
    except MissingDependency as exc:
        return fail(str(exc), action=exc.action)
    except Exception as exc:
        return fail(str(exc))


def extract_text_file(file_path, request):
    path = Path(file_path)
    if path.suffix.lower() == ".csv":
        rows = []
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            reader = csv.reader(handle)
            for row in reader:
                rows.append(" | ".join(cell.strip() for cell in row))
        return [unit("\n".join(rows), request, {"format": "csv"})]
    return [unit(read_text(path), request, {"format": path.suffix.lower().lstrip(".") or "text"})]


def extract_glossary_rows(request):
    try:
        file_path = request.get("filePath") or ""
        text = request.get("text") or ""
        relative_path = request.get("relativePath") or request.get("fileName") or ""
        title = request.get("title") or request.get("fileName") or "source"
        suffix = Path(file_path or relative_path or title).suffix.lower()

        if suffix == ".csv":
            rows = glossary_rows_from_csv_text(text) if text else glossary_rows_from_csv_path(file_path)
            return {"status": "ok", "rows": rows, "warnings": []}
        if suffix in {".xlsx", ".xlsm"}:
            return {"status": "ok", "rows": glossary_rows_from_xlsx(file_path), "warnings": []}
        return {"status": "ok", "rows": [], "warnings": []}
    except MissingDependency as exc:
        return fail(str(exc), action=exc.action)
    except Exception as exc:
        return fail(str(exc))


def glossary_rows_from_csv_path(file_path):
    with Path(file_path).open("r", encoding="utf-8-sig", newline="") as handle:
        return glossary_rows_from_csv_reader(csv.reader(handle), {})


def glossary_rows_from_csv_text(text):
    return glossary_rows_from_csv_reader(csv.reader(io.StringIO(str(text or ""))), {})


def glossary_rows_from_csv_reader(reader, locator_base):
    rows = list(reader)
    if not rows:
        return []
    headers = [str(cell or "").strip() for cell in rows[0]]
    out = []
    for row_index, row in enumerate(rows[1:], start=2):
        cells = {headers[i] or f"column_{i + 1}": str(row[i] if i < len(row) else "").strip() for i in range(len(headers))}
        if any(cells.values()):
            out.append({"cells": cells, "locator": {**locator_base, "rowRange": str(row_index)}})
    return out


def glossary_rows_from_xlsx(file_path):
    require_module("openpyxl", "Install openpyxl with `npm run py:deps`.")
    import openpyxl

    workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    out = []
    for sheet in workbook.worksheets:
        iterator = sheet.iter_rows(values_only=True)
        headers = None
        for row_index, row in enumerate(iterator, start=1):
            values = [format_cell(value) for value in row]
            if headers is None:
                if any(values):
                    headers = [value.strip() or f"column_{i + 1}" for i, value in enumerate(values)]
                continue
            cells = {headers[i] if i < len(headers) else f"column_{i + 1}": values[i] if i < len(values) else "" for i in range(max(len(headers), len(values)))}
            if any(cells.values()):
                out.append({"cells": cells, "locator": {"sheet": sheet.title, "rowRange": str(row_index)}})
    return out


def extract_pdf(file_path, request):
    require_module("fitz", "Install PyMuPDF with `npm run py:deps`.")
    import fitz

    docs = []
    warnings = []
    pdf = fitz.open(file_path)
    for page_index in range(pdf.page_count):
        page = pdf.load_page(page_index)
        text = page.get_text("text").strip()
        if len(text) < 20:
            ocr_text = ocr_pdf_page(page, request)
            if ocr_text:
                text = ocr_text
            else:
                warnings.append(f"page {page_index + 1}: no selectable text and OCR unavailable")
        if text:
            docs.append(unit(text, request, {"format": "pdf"}, {"page": page_index + 1}))
    if warnings and docs:
        docs[0]["metadata"]["warnings"] = warnings
    return docs


def ocr_pdf_page(page, request):
    if not shutil.which("tesseract"):
        return ""
    require_module("PIL", "Install Pillow with `npm run py:deps`.")
    require_module("pytesseract", "Install pytesseract with `npm run py:deps`.")
    import fitz
    from PIL import Image
    import pytesseract

    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
    mode = "RGB" if pix.alpha == 0 else "RGBA"
    image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
    return pytesseract.image_to_string(image, lang=request.get("ocrLanguages") or "kor+eng").strip()


def extract_docx(file_path, request):
    require_module("docx", "Install python-docx with `npm run py:deps`.")
    import docx

    document = docx.Document(file_path)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                paragraphs.append(" | ".join(cells))
    return [unit("\n".join(paragraphs), request, {"format": "docx"})]


def extract_pptx(file_path, request):
    require_module("pptx", "Install python-pptx with `npm run py:deps`.")
    from pptx import Presentation

    prs = Presentation(file_path)
    docs = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        if texts:
            docs.append(unit("\n".join(texts), request, {"format": "pptx"}, {"slide": index}))
    return docs


def extract_xlsx(file_path, request):
    require_module("openpyxl", "Install openpyxl with `npm run py:deps`.")
    import openpyxl

    workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    docs = []
    for sheet in workbook.worksheets:
        rows = []
        row_start = None
        row_end = None
        for row in sheet.iter_rows():
            values = [format_cell(cell.value) for cell in row]
            if any(values):
                if row_start is None:
                    row_start = row[0].row
                row_end = row[0].row
                rows.append(" | ".join(values))
            if len(rows) >= 80:
                docs.append(unit("\n".join(rows), request, {"format": "xlsx"}, {"sheet": sheet.title, "rowRange": f"{row_start}-{row_end}"}))
                rows = []
                row_start = None
                row_end = None
        if rows:
            docs.append(unit("\n".join(rows), request, {"format": "xlsx"}, {"sheet": sheet.title, "rowRange": f"{row_start}-{row_end}"}))
    return docs


def extract_image(file_path, request):
    if not shutil.which("tesseract"):
        raise MissingDependency("Tesseract binary is required for image OCR.", "Install tesseract and rerun indexing.")
    require_module("PIL", "Install Pillow with `npm run py:deps`.")
    require_module("pytesseract", "Install pytesseract with `npm run py:deps`.")
    from PIL import Image
    import pytesseract

    image = Image.open(file_path)
    text = pytesseract.image_to_string(image, lang=request.get("ocrLanguages") or "kor+eng").strip()
    return [unit(text, request, {"format": "image-ocr"})] if text else []


def extract_legacy_office(file_path, request, target_suffix):
    soffice = shutil.which("soffice")
    if not soffice:
        raise MissingDependency("LibreOffice `soffice` is required for legacy Office files.", "Install LibreOffice or convert this file first.")
    with tempfile.TemporaryDirectory() as temp_dir:
        subprocess.run(
            [soffice, "--headless", "--convert-to", target_suffix.lstrip("."), "--outdir", temp_dir, file_path],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        converted = next(Path(temp_dir).glob(f"*{target_suffix}"), None)
        if not converted:
            raise RuntimeError("LibreOffice conversion did not produce an output file.")
        next_request = dict(request)
        next_request["filePath"] = str(converted)
        if target_suffix == ".docx":
            return extract_docx(str(converted), next_request)
        if target_suffix == ".pptx":
            return extract_pptx(str(converted), next_request)
        if target_suffix == ".xlsx":
            return extract_xlsx(str(converted), next_request)
        raise RuntimeError(f"Unsupported conversion target {target_suffix}")


def unit(text, request, metadata=None, locator=None):
    locator = dict(locator or {})
    if request.get("relativePath"):
        locator["relativePath"] = request["relativePath"]
    return {
        "text": text or "",
        "locator": locator,
        "metadata": {
            "title": request.get("title") or request.get("fileName") or "source",
            "sourceId": request.get("sourceId") or "",
            **(metadata or {}),
        },
    }


def ok(documents, warnings=None):
    clean = [doc for doc in documents if (doc.get("text") or "").strip()]
    if not clean:
        return fail("No indexable text extracted.")
    return {"status": "ok", "documents": clean, "warnings": warnings or []}


def fail(error, action="Check the source file and rerun indexing."):
    return {"status": "failed_with_action", "error": error, "action": action, "documents": [], "warnings": []}


def has_module(name):
    return importlib.util.find_spec(name) is not None


def require_module(name, action):
    if not has_module(name):
        raise MissingDependency(f"Missing Python module `{name}`.", action)


def read_text(path):
    for encoding in ("utf-8-sig", "utf-8", "cp949", "latin-1"):
        try:
            return Path(path).read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return Path(path).read_text(errors="ignore")


def format_cell(value):
    if value is None:
        return ""
    return str(value).replace("\n", " ").strip()


class MissingDependency(Exception):
    def __init__(self, message, action):
        super().__init__(message)
        self.action = action


if __name__ == "__main__":
    main()
